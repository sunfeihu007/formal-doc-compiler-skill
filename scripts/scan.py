#!/usr/bin/env python3
"""Scan documents for terms in a compliance wordlist.

Usage:
    python3 scan.py --doc path/to/draft.docx --wordlist .compliance/wordlist.yaml
    python3 scan.py --doc a.docx --doc b.pdf --wordlist .compliance/wordlist.yaml

Wordlist semantics:
    - Terms are literal by default. Literal matching is case-insensitive and
      NFKC-normalized (full-width GLM variants match half-width GLM).
    - A term starting with "regex:" is a regular expression (also matched
      case-insensitively against NFKC-normalized text).
    - Invalid regex terms are reported as warning lines, not crashes.

Output: one JSON object per line — hits, then warnings, then a summary line
with {"summary": true, "hits": N, "warnings": N, "terms_checked": N, "docs": [...]}.
Exit code: 0 on a completed scan (with or without hits — read the summary
line), 1 on fatal errors (missing file, unsupported format).
"""
import argparse
import json
import re
import sys
import unicodedata
from pathlib import Path

REGEX_PREFIX = "regex:"


def extract_text(path: Path) -> str:
    suffix = path.suffix.lower()
    if suffix == ".docx":
        from docx import Document
        d = Document(str(path))
        parts = [p.text for p in d.paragraphs]
        for t in d.tables:
            for row in t.rows:
                for cell in row.cells:
                    parts.append(cell.text)
        # headers/footers are easy to forget and often hold boilerplate
        for section in d.sections:
            for hf in (section.header, section.footer):
                for p in hf.paragraphs:
                    parts.append(p.text)
        return "\n".join(parts)
    if suffix == ".pdf":
        import subprocess
        return subprocess.check_output(
            ["pdftotext", "-layout", str(path), "-"], text=True
        )
    if suffix == ".pptx":
        from pptx import Presentation
        p = Presentation(str(path))
        parts = []
        for slide in p.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    parts.append(shape.text_frame.text)
            if slide.has_notes_slide:
                parts.append(slide.notes_slide.notes_text_frame.text)
        return "\n".join(parts)
    if suffix == ".xlsx":
        import openpyxl
        wb = openpyxl.load_workbook(str(path), data_only=True, read_only=True)
        parts = []
        for ws in wb.worksheets:
            for row in ws.iter_rows(values_only=True):
                for v in row:
                    if v is not None:
                        parts.append(str(v))
        return "\n".join(parts)
    if suffix in {".md", ".txt", ".html"}:
        return path.read_text(encoding="utf-8")
    raise ValueError(f"unsupported extension: {suffix} "
                     "(supported: .docx .pdf .pptx .xlsx .md .txt .html)")


def compile_term(term: str):
    """Return a compiled pattern for a wordlist term, or raise re.error."""
    if term.startswith(REGEX_PREFIX):
        pattern = term[len(REGEX_PREFIX):]
    else:
        pattern = re.escape(unicodedata.normalize("NFKC", term))
    return re.compile(pattern, re.IGNORECASE)


def scan(text: str, wordlist: dict):
    """Scan NFKC-normalized text; return (hits, warnings)."""
    norm = unicodedata.normalize("NFKC", text)
    hits, warnings = [], []
    for cat, spec in (wordlist.get("categories") or {}).items():
        spec = spec or {}
        suggestion = spec.get("suggested_replacement", "")
        for term in spec.get("terms") or []:
            term = str(term)
            try:
                pattern = compile_term(term)
            except re.error as e:
                warnings.append({
                    "warning": True,
                    "category": cat,
                    "term": term,
                    "error": str(e),
                })
                continue
            for m in pattern.finditer(norm):
                start = max(0, m.start() - 30)
                end = min(len(norm), m.end() + 30)
                ctx = norm[start:end].replace("\n", " ")
                hits.append({
                    "category": cat,
                    "term": term,
                    "matched": m.group(0),
                    "context": ctx,
                    "suggested_replacement": suggestion,
                })
    return hits, warnings


def count_terms(wordlist: dict) -> int:
    return sum(
        len(spec.get("terms") or [])
        for spec in (wordlist.get("categories") or {}).values()
        if spec
    )


def main():
    ap = argparse.ArgumentParser(
        description="Scan documents for forbidden terms from a YAML wordlist.")
    ap.add_argument("--doc", required=True, action="append",
                    help="document to scan (repeatable)")
    ap.add_argument("--wordlist", required=True)
    args = ap.parse_args()

    wl_path = Path(args.wordlist)
    if not wl_path.exists():
        print(json.dumps({"error": f"wordlist not found: {wl_path}"}))
        sys.exit(1)

    import yaml
    wordlist = yaml.safe_load(wl_path.read_text(encoding="utf-8")) or {}

    all_hits, all_warnings = [], []
    for doc in args.doc:
        doc_path = Path(doc)
        if not doc_path.exists():
            print(json.dumps({"error": f"document not found: {doc_path}"}))
            sys.exit(1)
        try:
            text = extract_text(doc_path)
        except ValueError as e:
            print(json.dumps({"error": f"{doc_path}: {e}"}))
            sys.exit(1)
        hits, warnings = scan(text, wordlist)
        for h in hits:
            h["doc"] = str(doc_path)
        all_hits.extend(hits)
        all_warnings.extend(warnings)

    # warnings are deduped across docs (same bad term reported once)
    seen = set()
    unique_warnings = []
    for w in all_warnings:
        key = (w["category"], w["term"])
        if key not in seen:
            seen.add(key)
            unique_warnings.append(w)

    for h in all_hits:
        print(json.dumps(h, ensure_ascii=False))
    for w in unique_warnings:
        print(json.dumps(w, ensure_ascii=False))

    print(json.dumps({
        "summary": True,
        "hits": len(all_hits),
        "warnings": len(unique_warnings),
        "terms_checked": count_terms(wordlist),
        "docs": args.doc,
    }, ensure_ascii=False))


if __name__ == "__main__":
    main()
